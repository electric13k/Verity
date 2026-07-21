"""G9 file-output tools — produce deliverables from markdown.

Four format tools (docx / pptx / xlsx / pdf) turn model-authored markdown into a
real office/pdf file, plus a gated ``generate_image`` stub. Each tool registers
into the same G1 registry and returns a prompt-safe reference to the stored
deliverable (never the bytes) through the ``prompt_safe`` choke point.

Degrade, never die: the generator library is an OPTIONAL extra, imported LAZILY
inside ``run``; if it is not installed the tool returns a clear
"<lib> not installed" result instead of raising. ``generate_image`` degrades to
"image generation not configured" until an image provider key is present.

Tenant law: every tool stores under the caller's tenant directory
(``filestore``, keyed by ``tenant.user_id`` from gRPC metadata only).
"""

from __future__ import annotations

import re
from abc import abstractmethod

from app.config import settings
from app.tenant import TenantCtx
from app.tools import filestore
from app.tools.base import Tool, ToolResult, prompt_safe

_MD_TABLE_ROW = re.compile(r"^\s*\|(.+)\|\s*$")


def _ref_result(tool: str, deliverable: filestore.Deliverable) -> ToolResult:
    body = (
        f"Created {tool} deliverable.\n"
        f"file_id: {deliverable.file_id}\n"
        f"name: {deliverable.name}\n"
        f"content_type: {deliverable.content_type}\n"
        f"bytes: {deliverable.byte_size}\n"
        "The file is stored; reference it by file_id."
    )
    return ToolResult(prompt_safe(body, source=f"{tool}:{deliverable.file_id}"))


def _parse_table(markdown: str) -> list[list[str]]:
    """Rows of the first markdown table found, else []. The separator row
    (---|---) is dropped."""
    rows: list[list[str]] = []
    for line in markdown.splitlines():
        m = _MD_TABLE_ROW.match(line)
        if not m:
            if rows:
                break
            continue
        cells = [c.strip() for c in m.group(1).split("|")]
        if all(set(c) <= {"-", ":", " "} and c for c in cells):
            continue  # separator row
        rows.append(cells)
    return rows


class _MarkdownDeliverableTool(Tool):
    """Shared shape: (markdown, filename?) → stored deliverable reference."""

    ext = ""
    lib_import = ""  # module name to lazily import; "" = no external lib

    def __init__(self) -> None:
        self.parameters = {
            "type": "object",
            "properties": {
                "markdown": {
                    "type": "string",
                    "description": "The document content as markdown.",
                },
                "filename": {
                    "type": "string",
                    "description": "Base name for the file (extension is added).",
                },
            },
            "required": ["markdown"],
        }

    async def run(self, args: dict, tenant: TenantCtx) -> ToolResult:
        markdown = args.get("markdown")
        if not isinstance(markdown, str) or not markdown.strip():
            return ToolResult(
                prompt_safe("markdown content is required", source=self.name),
                is_error=True,
            )
        filename = (args.get("filename") or "deliverable").strip() or "deliverable"
        try:
            data = self._render(markdown)
        except ImportError:
            return ToolResult(
                prompt_safe(
                    f"{self.name} unavailable: the '{self.lib_import}' library is "
                    "not installed on this server (optional extra)",
                    source=self.name,
                ),
            )
        except Exception as exc:  # a bad render must not break the loop
            return ToolResult(
                prompt_safe(f"could not build {self.ext}: {exc}", source=self.name),
                is_error=True,
            )
        deliverable = filestore.store(tenant.user_id, filename, self.ext, data)
        return _ref_result(self.name, deliverable)

    @abstractmethod
    def _render(self, markdown: str) -> bytes:
        """Render markdown → file bytes. Raises ImportError if the lib is absent."""


class DocxTool(_MarkdownDeliverableTool):
    name = "create_docx"
    description = "Create a Microsoft Word (.docx) document from markdown."
    ext = "docx"
    lib_import = "python-docx"

    def _render(self, markdown: str) -> bytes:
        import io

        from docx import Document  # lazy: optional extra

        doc = Document()
        for raw in markdown.splitlines():
            line = raw.rstrip()
            if not line.strip():
                continue
            heading = re.match(r"^(#{1,6})\s+(.*)$", line)
            bullet = re.match(r"^\s*[-*+]\s+(.*)$", line)
            if heading:
                doc.add_heading(heading.group(2), level=min(len(heading.group(1)), 4))
            elif bullet:
                doc.add_paragraph(bullet.group(1), style="List Bullet")
            else:
                doc.add_paragraph(line)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


class PptxTool(_MarkdownDeliverableTool):
    name = "create_pptx"
    description = (
        "Create a PowerPoint (.pptx) deck from markdown. Each top-level heading "
        "starts a new slide; bullet lines become slide content."
    )
    ext = "pptx"
    lib_import = "python-pptx"

    def _render(self, markdown: str) -> bytes:
        import io

        from pptx import Presentation  # lazy: optional extra

        prs = Presentation()
        title_layout = prs.slide_layouts[1]  # title + content

        slides: list[tuple[str, list[str]]] = []
        current: tuple[str, list[str]] | None = None
        for raw in markdown.splitlines():
            line = raw.rstrip()
            if not line.strip():
                continue
            heading = re.match(r"^(#{1,3})\s+(.*)$", line)
            bullet = re.match(r"^\s*[-*+]\s+(.*)$", line)
            if heading:
                if current:
                    slides.append(current)
                current = (heading.group(2), [])
            else:
                if current is None:
                    current = ("Slide", [])
                current[1].append(bullet.group(1) if bullet else line)
        if current:
            slides.append(current)
        if not slides:
            slides = [("Deliverable", [markdown.strip()[:400]])]

        for title, bullets in slides:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title[:200]
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(bullets[:20]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = b[:400]
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


class XlsxTool(_MarkdownDeliverableTool):
    name = "create_xlsx"
    description = (
        "Create an Excel (.xlsx) spreadsheet from markdown. A markdown table "
        "becomes rows/columns; otherwise each line becomes a row."
    )
    ext = "xlsx"
    lib_import = "openpyxl"

    def _render(self, markdown: str) -> bytes:
        import io

        from openpyxl import Workbook  # lazy: optional extra

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        table = _parse_table(markdown)
        if table:
            for row in table:
                ws.append(row)
        else:
            for line in markdown.splitlines():
                if line.strip():
                    ws.append([line])
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()


class PdfTool(_MarkdownDeliverableTool):
    name = "create_pdf"
    description = "Create a PDF document from markdown."
    ext = "pdf"
    lib_import = "reportlab"

    def _render(self, markdown: str) -> bytes:
        import io

        from reportlab.lib.pagesizes import LETTER  # lazy: optional extra
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
        from xml.sax.saxutils import escape

        styles = getSampleStyleSheet()
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=LETTER)
        flow = []
        for raw in markdown.splitlines():
            line = raw.rstrip()
            if not line.strip():
                flow.append(Spacer(1, 6))
                continue
            heading = re.match(r"^(#{1,6})\s+(.*)$", line)
            bullet = re.match(r"^\s*[-*+]\s+(.*)$", line)
            if heading:
                level = min(len(heading.group(1)), 4)
                flow.append(Paragraph(escape(heading.group(2)), styles[f"Heading{level}"]))
            elif bullet:
                flow.append(Paragraph("• " + escape(bullet.group(1)), styles["Normal"]))
            else:
                flow.append(Paragraph(escape(line), styles["Normal"]))
        if not flow:
            flow.append(Paragraph("(empty)", styles["Normal"]))
        doc.build(flow)
        return buf.getvalue()


class ImageGenTool(Tool):
    """Gated image generation. Wired for a provider key; degrades cleanly."""

    name = "generate_image"
    description = (
        "Generate an image from a text prompt. Returns a stored image reference. "
        "Requires an image-generation provider to be configured on the server."
    )
    parameters = {
        "type": "object",
        "properties": {
            "prompt": {"type": "string", "description": "Description of the image."},
        },
        "required": ["prompt"],
    }

    async def run(self, args: dict, tenant: TenantCtx) -> ToolResult:
        prompt = (args.get("prompt") or "").strip()
        if not prompt:
            return ToolResult(
                prompt_safe("an image prompt is required", source=self.name),
                is_error=True,
            )
        if not settings.image_api_key:
            return ToolResult(
                prompt_safe(
                    "image generation is not configured on this server; no image "
                    "was produced",
                    source=self.name,
                ),
            )
        # Key present but no provider adapter is wired yet: degrade gracefully
        # rather than pretend. (Provider integration lands when a key + provider
        # choice are supplied — never fabricate an image.)
        return ToolResult(
            prompt_safe(
                "image generation is enabled but no provider adapter is wired on "
                "this build; no image was produced",
                source=self.name,
            ),
        )


def file_output_tools() -> list[Tool]:
    return [DocxTool(), PptxTool(), XlsxTool(), PdfTool(), ImageGenTool()]
